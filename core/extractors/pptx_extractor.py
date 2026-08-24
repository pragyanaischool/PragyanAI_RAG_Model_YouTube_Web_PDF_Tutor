"""
PPTX Extractor Module for NCET GenAI Multimodal RAG.
Handles Microsoft PowerPoint (.pptx) file streams, extracting slide titles, 
body text, shape paragraphs, embedded slide tables, and speaker/presenter notes per slide.
"""

import io
import re
from typing import Any, Dict, Optional
from pptx import Presentation


def clean_slide_text(text: str) -> str:
    """
    Normalizes whitespace and removes excessive newline artifacts.
    """
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def extract_from_pptx(uploaded_file: Any) -> Optional[Dict[str, Any]]:
    """
    Parses an uploaded Microsoft PowerPoint (.pptx) file stream or file-like binary object,
    extracting text elements slide by slide, preserving structure, and capturing speaker notes.

    Args:
        uploaded_file: Streamlit UploadedFile, file-like object, or bytes buffer.

    Returns:
        Optional[Dict[str, Any]]: Document payload containing:
            - 'source': Presentation filename
            - 'type': 'pptx'
            - 'title': Presentation filename or title
            - 'slide_count': Total number of slides
            - 'content': Structured text payload with per-slide delimiters
        Returns None if extraction fails or file contains no extractable text.
    """
    if uploaded_file is None:
        return None

    file_name = getattr(uploaded_file, "name", "Presentation.pptx")

    try:
        # Support Streamlit UploadedFile, BytesIO, or raw byte streams
        if hasattr(uploaded_file, "getvalue"):
            pptx_stream = io.BytesIO(uploaded_file.getvalue())
        elif isinstance(uploaded_file, bytes):
            pptx_stream = io.BytesIO(uploaded_file)
        else:
            pptx_stream = uploaded_file

        prs = Presentation(pptx_stream)
        total_slides = len(prs.slides)

        if total_slides == 0:
            print(f"[Warning in pptx_extractor]: {file_name} contains zero slides.")
            return None

        slides_text = []
        valid_slide_count = 0

        for idx, slide in enumerate(prs.slides):
            slide_lines = [f"--- Slide {idx + 1} of {total_slides} ---"]
            slide_has_content = False

            # 1. Extract Slide Title
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = clean_slide_text(slide.shapes.title.text)
                slide_lines.append(f"Title: {title_text}")
                slide_has_content = True

            # 2. Extract Text from all shapes and text frames
            for shape in slide.shapes:
                # Skip title shape as it's already processed above
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        clean_para = clean_slide_text(para.text)
                        if clean_para:
                            # Preserve bullet level indentation if present
                            level_indent = "  " * getattr(para, "level", 0)
                            bullet_prefix = f"{level_indent}- " if getattr(para, "level", 0) > 0 else ""
                            slide_lines.append(f"{bullet_prefix}{clean_para}")
                            slide_has_content = True

                # 3. Extract Slide Tables
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        row_cells = [clean_slide_text(cell.text) for cell in row.cells]
                        non_empty_cells = [c for c in row_cells if c]
                        if non_empty_cells:
                            table_rows.append(" | ".join(non_empty_cells))
                    if table_rows:
                        slide_lines.append("[Slide Table]:\n" + "\n".join(table_rows))
                        slide_has_content = True

            # 4. Extract Speaker Notes (if present)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                speaker_notes = clean_slide_text(slide.notes_slide.notes_text_frame.text)
                if speaker_notes:
                    slide_lines.append(f"[Speaker Notes]: {speaker_notes}")
                    slide_has_content = True

            if slide_has_content:
                slides_text.append("\n".join(slide_lines))
                valid_slide_count += 1

        if not slides_text:
            print(f"[Warning in pptx_extractor]: No readable text found in {file_name}.")
            return None

        full_presentation_text = "\n\n".join(slides_text)

        header_block = (
            f"[Source: PPTX Presentation]\n"
            f"[Title: {file_name}]\n"
            f"[Total Slides: {total_slides} | Content Slides: {valid_slide_count}]\n"
            f"----------------------------------------\n\n"
        )

        full_content = f"{header_block}{full_presentation_text}"

        return {
            "source": file_name,
            "type": "pptx",
            "title": file_name,
            "slide_count": total_slides,
            "content": full_content,
        }

    except Exception as e:
        print(f"[Error in pptx_extractor]: Failed parsing {file_name}: {e}")
        return None
